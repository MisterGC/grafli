"""Export a flow as an editable PowerPoint (.pptx) presentation.

Hybrid by design: the bookmark's framed diagram region is a flattened picture
(rasterized via ``QGraphicsScene.render``, exactly as the PDF exporter does), but
everything textual — title, progress, caption, footer and in-place notes — is a
native, editable PowerPoint textbox. So the deck is directly usable on export yet
fully tweakable: retext, restyle, add or remove slides.

Theme-friendly: text is styled through *theme* colour and font references rather
than hardcoded values, so a global theme/font swap in PowerPoint (Design ▸
Variants, or applying a corporate ``.thmx``) cascades onto it. Two presets:

- ``grafli`` (default) — injects grafli's palette (paper background, ``#D4804E``
  accent) and font into the deck theme, plus the title accent rule and bar
  dividers, so the export looks like grafli out of the box.
- ``blank`` — leaves the neutral default Office theme and drops the decorative
  accent rule/dividers: the cleanest base for merging into a corporate template.

The slide-typing/decision layer is shared with the PDF exporter via
``grafli.slideplan.build_slide_plan`` so both formats stay in lock-step.
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Pt
from PySide6.QtCore import QBuffer, QByteArray, QIODevice, QRectF, Qt
from PySide6.QtGui import QFont, QImage, QPainter, QTextDocument

from grafli.constants import FONT_FAMILY, resolve_textsize_px
from grafli.flows import isolate_focus, render_thumbnail_art
from grafli.md_note import is_md_note, md_body
from grafli.slideplan import SlidePlan, build_slide_plan, playback_text_fit

# PowerPoint 16:9 canvas in points — same fixed geometry as the PDF exporter, so
# the fractional layout math below mirrors pdfexport on a 960x540 pt page. When
# exporting onto a template we adopt *its* page size instead (read off the deck),
# so these are only the from-scratch defaults; all layout math reads the live
# page size from the export context (``_Ctx``).
_DEFAULT_PAGE_W = 960.0
_DEFAULT_PAGE_H = 540.0
_FOOTER_RESERVE_RATIO = 0.10
# python-pptx geometry is in EMU; 1 point = 12700 EMU. Used to read a template
# layout's placeholder rects (EMU) back into the point space the layout math uses.
_EMU_PER_PT = 12700.0

# Single-line text sizes in points.
_TITLE_PT = 27.0          # slide title bar
_PROGRESS_PT = 18.0       # "i / n"
_FOOTER_PT = 12.0
_COVER_TITLE_PT = 46.0

# Area-filling text uses a (min, ideal, max) point band and a grow-to-fill rule,
# mirroring the PDF exporter: PowerPoint's autofit only *shrinks* text, never
# grows it, so we compute the fill size ourselves and set it explicitly. Bands
# match grafli.pdfexport so the PPTX text fills the slide exactly like the PDF.
_BODY_BAND = (18.0, 24.0, 30.0)     # text-slide hero (fallback for dense notes)
_DESC_BAND = (14.0, 18.0, 22.0)     # cover description
# A text slide normally sizes at playback parity (the note zoomed to fill the
# hero, like fitInView frames it in the app) capped here; below the body-band
# minimum it falls back to the band's shrink-to-fit instead.
_BODY_MAX_PT = 60.0
_CAPTION_BAND = (12.0, 15.0, 18.0)
_FILL_FLOOR = 0.45                  # grow until the text fills this much height
_HEADING_SCALE = {1: 1.4, 2: 1.2, 3: 1.05}
# Default PowerPoint textbox internal margins (0.1in x 0.05in) in points — the
# fill fit subtracts these so the chosen size accounts for the inset.
_TF_MARGIN_W = 7.2
_TF_MARGIN_H = 3.6

# An in-place note overlay below this point size is too small to read on a slide;
# we place it anyway (faithful position wins) but flag the slide as overloaded.
_READABLE_MIN_PT = 11.0

# grafli theme palette (mirrors grafli.constants), injected into the deck theme
# so theme-referenced runs/shapes default to the grafli look yet stay swappable.
_GRAFLI_COLORS = {
    "dk1": "2F3437", "lt1": "E8E4DD", "dk2": "4A4A4A", "lt2": "F5F2ED",
    "accent1": "D4804E", "accent2": "004578", "accent3": "0178D4",
    "accent4": "4EBF71", "accent5": "D4BA6A", "accent6": "B0A1CA",
    "hlink": "2B6CB0", "folHlink": "805AD5",
}


@dataclass
class _Theme:
    name: str
    inject_grafli: bool   # rewrite the deck theme with the grafli palette/font
    chrome: bool          # draw the title accent rule + bar dividers


_THEMES = {
    "grafli": _Theme("grafli", inject_grafli=True, chrome=True),
    "blank": _Theme("blank", inject_grafli=False, chrome=False),
    # Template mode reuses the neutral preset (no grafli palette/chrome); its
    # extra behaviour (placeholder-driven placement, no own footer/progress) is
    # gated on ``_Ctx.template`` rather than the theme itself.
    "template": _Theme("template", inject_grafli=False, chrome=False),
}


@dataclass
class _Ctx:
    """Per-export geometry + mode, threaded through the slide builders so the
    layout math is page-size agnostic and template-aware."""
    page_w: float
    page_h: float
    theme: _Theme
    footer: str
    template: bool = False

# Theme-font reference tokens — resolve to the deck's major/minor latin fonts, so
# a font-theme change in PowerPoint cascades onto the text.
_FONT_MAJOR = "+mj-lt"
_FONT_MINOR = "+mn-lt"


def export_flow_to_pptx(view, flow, out_path: str | Path,
                        theme: str = "grafli", template: str | Path | None = None,
                        title_layout: str | None = None,
                        content_layout: str | None = None) -> tuple[int, list]:
    """Render ``flow`` to a .pptx at ``out_path``.

    Two modes:
    - From scratch (``template`` is None): build a fresh deck at the default
      16:9 geometry and style it with the ``theme`` preset (``grafli`` /
      ``blank``).
    - Onto a template (``template`` is a .pptx path): open that deck as the base
      — keeping its master, layouts, theme and slide size — strip its slides, and
      drop grafli content onto fresh slides built from its layouts. ``title_layout``
      / ``content_layout`` name which layouts to use (the heuristic picks sensible
      defaults when omitted). The corporate theme cascades for free because every
      run already references theme fonts/colours.

    Returns ``(slide_count, overloaded)`` — same contract as the PDF exporter —
    where ``overloaded`` lists ``(step_index, title)`` for slides whose in-place
    notes fall below the readable floor, so the caller can warn the author.
    """
    board = view.board
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if template is not None:
        prs = Presentation(str(template))
        _strip_slides(prs)
        th = _THEMES["template"]
        ctx = _Ctx(prs.slide_width / _EMU_PER_PT, prs.slide_height / _EMU_PER_PT,
                   th, board.footer or "", template=True)
        tlayout = _resolve_layout(prs, title_layout, "title")
        clayout = _resolve_layout(prs, content_layout, "content")
    else:
        th = _THEMES.get(theme, _THEMES["grafli"])
        prs = Presentation()
        prs.slide_width = Pt(_DEFAULT_PAGE_W)
        prs.slide_height = Pt(_DEFAULT_PAGE_H)
        if th.inject_grafli:
            _apply_grafli_theme(prs)
        ctx = _Ctx(_DEFAULT_PAGE_W, _DEFAULT_PAGE_H, th, board.footer or "")
        tlayout = clayout = prs.slide_layouts[6]   # the built-in blank layout

    plans = build_slide_plan(view, flow)

    # Clear selection and suppress the paper background so rasterized regions are
    # transparent outside their items (they then blend onto the slide), restoring
    # both afterwards — mirrors the PDF exporter.
    sel = list(view._scene.selectedItems())
    for item in sel:
        item.setSelected(False)
    old_bg = view._scene.backgroundBrush()
    view._scene.setBackgroundBrush(Qt.GlobalColor.transparent)

    overloaded = []
    try:
        _build_title_slide(prs.slides.add_slide(tlayout), view, board, flow, ctx)
        for plan in plans[1:]:
            slide = prs.slides.add_slide(clayout)
            if _build_content_slide(slide, view, plan, ctx):
                overloaded.append((plan.index, plan.title))
    finally:
        view._scene.setBackgroundBrush(old_bg)
        for item in sel:
            item.setSelected(True)

    prs.save(str(out_path))
    return len(plans), overloaded


# ── template ──────────────────────────────────────────────────────────────────

def _strip_slides(prs) -> None:
    """Remove every slide from an opened template, leaving its masters, layouts
    and theme intact — so we add our own slides onto a clean corporate base."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rId = sld_id.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id)


def _all_layouts(prs) -> list:
    return [lo for m in prs.slide_masters for lo in m.slide_layouts]


def _is_title_ph(ph) -> bool:
    return "TITLE" in str(ph.placeholder_format.type)


def _ph_area(ph) -> float:
    return float((ph.width or 0)) * float((ph.height or 0))


def _body_area(layout) -> float:
    """Largest non-title placeholder area on a layout — its candidate body zone."""
    areas = [_ph_area(ph) for ph in layout.placeholders if not _is_title_ph(ph)]
    return max(areas, default=0.0)


def _resolve_layout(prs, name: str | None, kind: str):
    """The slide layout to use for ``kind`` ('title' | 'content').

    An exact ``name`` match wins. Otherwise heuristics: a title slide prefers a
    layout named '*title*' that has a title placeholder; a content slide prefers
    the layout with the largest body placeholder (the most room for the diagram).
    Falls back to the first layout."""
    layouts = _all_layouts(prs)
    if name:
        for lo in layouts:
            if lo.name == name:
                return lo
    if kind == "title":
        named = [lo for lo in layouts
                 if "title" in (lo.name or "").lower()
                 and any(_is_title_ph(ph) for ph in lo.placeholders)]
        if named:
            return named[0]
        with_title = [lo for lo in layouts
                      if any(_is_title_ph(ph) for ph in lo.placeholders)]
        if with_title:
            return with_title[0]
    else:
        best = max(layouts, key=_body_area, default=None)
        if best is not None and _body_area(best) > 0:
            return best
    return layouts[0] if layouts else prs.slide_layouts[0]


def _classify_placeholders(slide):
    """Split a freshly-added slide's placeholders into (title, body).

    Title = a real TITLE/CENTER_TITLE placeholder, else the top-most remaining
    one. Body = the largest non-title placeholder (where the diagram goes). Either
    may be None. When only one placeholder exists it is treated as the body."""
    phs = list(slide.placeholders)
    if not phs:
        return None, None
    titles = [p for p in phs if _is_title_ph(p)]
    body = max((p for p in phs if p not in titles), key=_ph_area, default=None)
    if body is None:
        body = max(phs, key=_ph_area)
    if titles:
        title = titles[0]
    else:
        rest = [p for p in phs if p is not body]
        title = min(rest, key=lambda p: (p.top or 0)) if rest else None
    return (title if title is not body else None), body


def _ph_rect(ph) -> QRectF:
    """A placeholder's rect in points (the space the layout math works in)."""
    return QRectF(float(ph.left) / _EMU_PER_PT, float(ph.top) / _EMU_PER_PT,
                  float(ph.width) / _EMU_PER_PT, float(ph.height) / _EMU_PER_PT)


def _set_ph_text(ph, text: str) -> None:
    """Set a placeholder's text while keeping its inherited (template) styling —
    so the title adopts the corporate font/size/colour rather than ours."""
    ph.text_frame.text = text


def _remove_ph(ph) -> None:
    """Drop an unused placeholder so it leaves no 'click to add text' prompt."""
    ph._element.getparent().remove(ph._element)


# ── theme ───────────────────────────────────────────────────────────────────

def _apply_grafli_theme(prs) -> None:
    """Rewrite the deck theme's colour scheme and latin fonts with grafli's, so
    theme-referenced text/shapes default to the grafli look while a later theme
    swap still overrides them. Best-effort: on any structural surprise we leave
    the default Office theme (text stays theme-referenced, just Office-coloured)."""
    try:
        part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        theme = etree.fromstring(part.blob)
        clr = theme.find(".//" + qn("a:clrScheme"))
        for slot, hexv in _GRAFLI_COLORS.items():
            el = clr.find(qn("a:" + slot))
            if el is None:
                continue
            for child in list(el):
                el.remove(child)
            etree.SubElement(el, qn("a:srgbClr"), val=hexv)
        fonts = theme.find(".//" + qn("a:fontScheme"))
        for which in ("a:majorFont", "a:minorFont"):
            latin = fonts.find(qn(which) + "/" + qn("a:latin"))
            if latin is not None:
                latin.set("typeface", FONT_FAMILY)
        part._blob = etree.tostring(theme, xml_declaration=True,
                                    encoding="UTF-8", standalone=True)
    except Exception:
        pass


# ── slides ──────────────────────────────────────────────────────────────────

def _build_title_slide(slide, view, board, flow, ctx: _Ctx) -> None:
    if ctx.template:
        _build_title_slide_template(slide, flow, ctx)
        return
    page_w, page_h, th = ctx.page_w, ctx.page_h, ctx.theme
    margin = page_h * 0.10
    if th.inject_grafli:
        _fill_background(slide, MSO_THEME_COLOR.BACKGROUND_1)
    if board is not None and board.title_bg == "thumbnail-art":
        art = render_thumbnail_art(view, board, flow, 2400,
                                   int(2400 * page_h / page_w))
        if art is not None:
            slide.shapes.add_picture(_png_stream(art), 0, 0,
                                     Pt(page_w), Pt(page_h))

    y = page_h * 0.26
    w = page_w - margin * 2
    box = slide.shapes.add_textbox(Pt(margin), Pt(y), Pt(w), Pt(page_h * 0.18))
    _fill_text(box.text_frame, flow.label, is_md=False, base_pt=_COVER_TITLE_PT,
               font_ref=_FONT_MAJOR, color=MSO_THEME_COLOR.TEXT_1, bold=True)

    ry = y + page_h * 0.20
    if th.chrome:
        _accent_rule(slide, margin, ry, w * 0.5, page_h * 0.006)

    if flow.description:
        drect = (margin, ry + page_h * 0.03, w, page_h * 0.55)
        box = slide.shapes.add_textbox(*(Pt(v) for v in drect))
        base = _fit_body_pt(flow.description, True, w - _TF_MARGIN_W * 2,
                            page_h * 0.55 - _TF_MARGIN_H * 2, _DESC_BAND)
        _fill_text(box.text_frame, flow.description, is_md=True,
                   base_pt=base, font_ref=_FONT_MINOR,
                   color=MSO_THEME_COLOR.TEXT_1, autosize=False)


def _build_title_slide_template(slide, flow, ctx: _Ctx) -> None:
    """Cover slide on a template: the flow label fills the template's title
    placeholder (inheriting its styling); the description, when present, goes in a
    textbox just below it. Other placeholders (date/footer/number) are left to the
    template."""
    phs = list(slide.placeholders)
    titles = [p for p in phs if _is_title_ph(p)]
    title_ph = titles[0] if titles else (max(phs, key=_ph_area) if phs else None)
    if title_ph is not None:
        _set_ph_text(title_ph, flow.label)
    if not flow.description:
        return
    if title_ph is not None:
        tr = _ph_rect(title_ph)
        dx, dw = tr.left(), tr.width()
        dy = tr.bottom() + tr.height() * 0.2
    else:
        margin = ctx.page_w * 0.08
        dx, dw, dy = margin, ctx.page_w - margin * 2, ctx.page_h * 0.45
    dh = max(40.0, dw * 0.30)
    box = slide.shapes.add_textbox(Pt(dx), Pt(dy), Pt(dw), Pt(dh))
    base = _fit_body_pt(flow.description, True, dw - _TF_MARGIN_W * 2,
                        dh - _TF_MARGIN_H * 2, _DESC_BAND)
    _fill_text(box.text_frame, flow.description, is_md=True, base_pt=base,
               font_ref=_FONT_MINOR, color=MSO_THEME_COLOR.TEXT_1, autosize=False)


def _build_content_slide(slide, view, plan: SlidePlan, ctx: _Ctx) -> bool:
    if ctx.template:
        return _build_content_slide_template(slide, view, plan, ctx)

    page_w, page_h, th, footer = ctx.page_w, ctx.page_h, ctx.theme, ctx.footer
    margin = page_h * 0.06
    if th.inject_grafli:
        _fill_background(slide, MSO_THEME_COLOR.BACKGROUND_1)
    if footer:
        _build_footer(slide, footer, ctx)
    has_title = bool(plan.title)

    hero_top = margin * 0.5
    hero_bottom = page_h - margin * 0.5 - _footer_reserve(ctx)

    if has_title:
        bar_h = page_h * 0.12
        tbox = slide.shapes.add_textbox(Pt(margin), Pt(margin * 0.5),
                                        Pt(page_w - margin * 2), Pt(bar_h))
        _fill_text(tbox.text_frame, plan.title, is_md=False, base_pt=_TITLE_PT,
                   font_ref=_FONT_MAJOR, color=MSO_THEME_COLOR.TEXT_1, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE)
        pbox = slide.shapes.add_textbox(Pt(margin), Pt(margin * 0.5),
                                        Pt(page_w - margin * 2), Pt(bar_h))
        _fill_text(pbox.text_frame, f"{plan.index + 1} / {plan.total}",
                   is_md=False, base_pt=_PROGRESS_PT, font_ref=_FONT_MINOR,
                   color=MSO_THEME_COLOR.TEXT_1, anchor=MSO_ANCHOR.MIDDLE,
                   align_right=True)
        rule_y = margin * 0.5 + bar_h
        if th.chrome:
            _divider(slide, margin, rule_y, page_w - margin * 2)
        hero_top = rule_y + margin * 0.5

    hero = QRectF(margin, hero_top, page_w - margin * 2, hero_bottom - hero_top)
    cap_bottom = page_h - margin - _footer_reserve(ctx)
    return _render_hero(slide, view, plan, hero, ctx, cap_bottom)


def _build_content_slide_template(slide, view, plan: SlidePlan,
                                  ctx: _Ctx) -> bool:
    """Content slide on a template: the title fills the template's title/heading
    placeholder, the diagram fits the body placeholder's region (so it respects
    the template's margins, header and footer), and grafli's own chrome, footer
    and progress counter are dropped — the template supplies its own."""
    title_ph, body_ph = _classify_placeholders(slide)
    if plan.title and title_ph is not None:
        _set_ph_text(title_ph, plan.title)
    elif title_ph is not None:
        _remove_ph(title_ph)

    if body_ph is not None:
        hero = _ph_rect(body_ph)
        _remove_ph(body_ph)            # we draw over it; drop the empty prompt
    else:
        margin = ctx.page_h * 0.06
        top = ctx.page_h * (0.18 if plan.title else 0.08)
        hero = QRectF(margin, top, ctx.page_w - margin * 2,
                      ctx.page_h - top - margin)
    return _render_hero(slide, view, plan, hero, ctx, hero.bottom())


def _render_hero(slide, view, plan: SlidePlan, hero: QRectF, ctx: _Ctx,
                 cap_bottom: float) -> bool:
    """Render a content plan into ``hero``: a text note as native centred text, a
    'no anchor' note, or the rasterized diagram with live note overlays. A caption
    card, when present, floats above ``cap_bottom``. Shared by both modes."""
    # Text slide: the single note as native, editable, centred text.
    if plan.kind == "text":
        return _build_text_hero(slide, hero, plan)

    source = plan.source
    if source is None:
        box = slide.shapes.add_textbox(Pt(hero.left()), Pt(hero.top()),
                                       Pt(hero.width()), Pt(hero.height()))
        _fill_text(box.text_frame, "no anchor to render", is_md=False,
                   base_pt=_DESC_BAND[1], font_ref=_FONT_MINOR,
                   color=MSO_THEME_COLOR.TEXT_1, anchor=MSO_ANCHOR.MIDDLE)
        return False

    # Fit the framed region into the hero, preserving aspect ratio.
    scale = min(hero.width() / source.width(), hero.height() / source.height())
    tw, th_ = source.width() * scale, source.height() * scale
    fitted = QRectF(hero.left() + (hero.width() - tw) / 2,
                    hero.top() + (hero.height() - th_) / 2, tw, th_)

    img = _render_region(view, plan, fitted)
    slide.shapes.add_picture(_png_stream(img), Pt(fitted.left()),
                             Pt(fitted.top()), Pt(fitted.width()),
                             Pt(fitted.height()))

    # Overlay each note as an editable textbox at its mapped position, sized to
    # the same scene scale the diagram was scaled by so it reads in place.
    overloaded = False
    for item in plan.overlays:
        if _build_note_overlay(slide, item, source, fitted, scale):
            overloaded = True

    if plan.caption:
        _build_caption(slide, plan.caption, ctx, hero, cap_bottom)
    return overloaded


def _build_text_hero(slide, hero: QRectF, plan: SlidePlan) -> bool:
    """The slide's single note as native text, sized at playback parity: the
    note block fitted into the hero like ``fitInView`` frames it in the app,
    its font scaled by the same factor (capped at ``_BODY_MAX_PT``). A note too
    dense for that falls back to the band's shrink-to-fit over the full hero."""
    note = plan.text_note
    is_md = is_md_note(note.text)
    body = md_body(note.text) if is_md else note.text
    fit = playback_text_fit(resolve_textsize_px(note.textsize, ""),
                            plan.text_rect, hero, _BODY_BAND[0], _BODY_MAX_PT)
    if fit is not None:
        base, rect = fit
    else:
        rect = hero
        base = _fit_body_pt(body, is_md, hero.width() - _TF_MARGIN_W * 2,
                            hero.height() - _TF_MARGIN_H * 2, _BODY_BAND)
    box = slide.shapes.add_textbox(Pt(rect.left()), Pt(rect.top()),
                                   Pt(rect.width()), Pt(rect.height()))
    _fill_text(box.text_frame, body, is_md=is_md, base_pt=base,
               font_ref=_FONT_MINOR, color=MSO_THEME_COLOR.TEXT_1,
               anchor=MSO_ANCHOR.MIDDLE, autosize=False)
    return False


def _build_note_overlay(slide, item, source: QRectF, fitted: QRectF,
                        scale: float) -> bool:
    note = item.note
    nr = item.sceneBoundingRect()
    # The PDF clips note overlays to the diagram region, so a note hanging mostly
    # outside the framed area renders as nothing there. A PowerPoint textbox can't
    # be clipped, so a mostly-outside note would otherwise show in full, jammed
    # against a slide edge. Approximate the clip: only place the note when it is
    # substantially inside the framed region.
    inter = nr.intersected(source)
    nr_area = (nr.width() * nr.height()) or 1.0
    if inter.isEmpty() or (inter.width() * inter.height()) < 0.55 * nr_area:
        return False
    is_md = is_md_note(note.text)
    body = md_body(note.text) if is_md else note.text
    mapped = QRectF(fitted.left() + (nr.left() - source.left()) * scale,
                    fitted.top() + (nr.top() - source.top()) * scale,
                    nr.width() * scale, nr.height() * scale)
    # The note's on-canvas font, scaled by the same factor as the diagram, in pt.
    scene_w = nr.width() or 1.0
    note_scale = mapped.width() / scene_w
    base_pt = max(1.0, resolve_textsize_px(note.textsize, "") * note_scale)
    pad = getattr(item, "_PAD", 0) * note_scale
    box = slide.shapes.add_textbox(
        Pt(mapped.left() + pad), Pt(mapped.top() + pad),
        Pt(max(1.0, mapped.width() - pad * 2)),
        Pt(max(1.0, mapped.height() - pad * 2)))
    _fill_text(box.text_frame, body, is_md=is_md, base_pt=base_pt,
               font_ref=_FONT_MINOR, color=MSO_THEME_COLOR.TEXT_1,
               autosize=False)
    return base_pt < _READABLE_MIN_PT


def _build_caption(slide, text: str, ctx: _Ctx, hero: QRectF,
                   cap_bottom: float) -> None:
    """A floating dark rounded card with light text near the bottom of the
    content area — matching the on-canvas/PDF playback caption. Card fill and text
    use theme colours (dark text-1 ground, light background-1 text) so they invert
    cleanly under any theme. Centred within ``hero``, its base at ``cap_bottom``."""
    from pptx.enum.shapes import MSO_SHAPE
    pad = ctx.page_h * 0.020
    card_w = min(ctx.page_w * 0.62, hero.width())
    card_h = ctx.page_h * 0.16
    card_x = hero.left() + (hero.width() - card_w) / 2
    card_y = cap_bottom - card_h
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(card_x),
                                  Pt(card_y), Pt(card_w), Pt(card_h))
    card.fill.solid()
    card.fill.fore_color.theme_color = MSO_THEME_COLOR.TEXT_1
    card.line.fill.background()
    card.shadow.inherit = False
    tf = card.text_frame
    tf.margin_left = tf.margin_right = Pt(pad)
    tf.margin_top = tf.margin_bottom = Pt(pad * 0.6)
    base = _fit_body_pt(text, True, card_w - pad * 2, card_h - pad * 1.2,
                        _CAPTION_BAND)
    _fill_text(tf, text, is_md=True, base_pt=base, font_ref=_FONT_MINOR,
               color=MSO_THEME_COLOR.BACKGROUND_1, anchor=MSO_ANCHOR.MIDDLE,
               clear=True, autosize=False)


def _build_footer(slide, footer: str, ctx: _Ctx) -> None:
    """The board-global branding line, muted and left-aligned at the bottom, with
    a thin rule above it (grafli theme only) — mirrors the PDF footer band."""
    page_w, page_h, th = ctx.page_w, ctx.page_h, ctx.theme
    margin = page_h * 0.06
    band_h = page_h * 0.05
    band_y = page_h - margin * 0.35 - band_h
    if th.chrome:
        _divider(slide, margin, band_y - page_h * 0.012, page_w - margin * 2)
    box = slide.shapes.add_textbox(Pt(margin), Pt(band_y),
                                   Pt(page_w - margin * 2), Pt(band_h))
    _fill_text(box.text_frame, footer, is_md=True, base_pt=_FOOTER_PT,
               font_ref=_FONT_MINOR, color=MSO_THEME_COLOR.TEXT_1,
               anchor=MSO_ANCHOR.MIDDLE)


def _accent_rule(slide, x: float, y: float, w: float, h: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y),
                                 Pt(w), Pt(max(1.5, h)))
    bar.fill.solid()
    bar.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    bar.line.fill.background()
    bar.shadow.inherit = False


def _divider(slide, x: float, y: float, w: float) -> None:
    """A thin neutral hairline under the title bar."""
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w),
                                  Pt(1.0))
    line.fill.solid()
    line.fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
    line.line.fill.background()
    line.shadow.inherit = False


def _fill_background(slide, theme_color) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.theme_color = theme_color


# ── text ─────────────────────────────────────────────────────────────────────

def _fit_body_pt(body: str, is_md: bool, w_pt: float, h_pt: float,
                 band, fill_floor: float = _FILL_FLOOR) -> float:
    """Largest point size from ``band`` (min, ideal, max) at which ``body``,
    reflowed to ``w_pt``, fits ``h_pt`` — grown toward max until it fills
    ``fill_floor`` of the height, shrunk toward min on overflow. Same grow-to-
    fill model as the PDF exporter, so the PPTX text fills like the PDF."""
    doc = QTextDocument()
    doc.setDocumentMargin(0)
    if is_md:
        doc.setMarkdown(body, QTextDocument.MarkdownFeature.MarkdownDialectGitHub)
    else:
        doc.setPlainText(body)
    f = QFont(FONT_FAMILY)

    def h_at(pt: float) -> float:
        f.setPointSizeF(pt)
        doc.setDefaultFont(f)
        doc.setTextWidth(max(1.0, w_pt))
        return doc.size().height()

    lo, ideal, hi = band
    size = ideal
    if h_at(size) > h_pt:
        while size > lo:
            size -= 1
            if h_at(size) <= h_pt:
                break
    else:
        while size < hi:
            if h_at(size + 1) > h_pt:
                break
            size += 1
            if h_at(size) >= fill_floor * h_pt:
                break
    return float(size)


def _fill_text(tf, body: str, *, is_md: bool, base_pt: float, font_ref: str,
               color, bold: bool = False, anchor=MSO_ANCHOR.TOP,
               align_right: bool = False, autosize: bool = True,
               clear: bool = True) -> None:
    """Render ``body`` into text frame ``tf`` as native runs.

    Markdown (when ``is_md``) is parsed once via ``QTextDocument`` — the same
    engine the canvas and PDF use — and translated into PowerPoint paragraphs and
    runs: bold/italic, inline code (monospace), links (real hyperlinks), headings
    (scaled + bold) and bullet/ordered lists (hand-prefixed markers, like the PDF
    draws). Runs reference the theme font/colour so a theme swap restyles them.
    """
    from pptx.enum.text import PP_ALIGN
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE if not autosize \
        else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    doc = QTextDocument()
    doc.setDocumentMargin(0)
    f = QFont(FONT_FAMILY)
    f.setPointSizeF(base_pt)
    doc.setDefaultFont(f)
    if is_md:
        doc.setMarkdown(body, QTextDocument.MarkdownFeature.MarkdownDialectGitHub)
    else:
        doc.setPlainText(body)

    if clear:
        tf.clear()
    first = True
    blk = doc.begin()
    while blk.isValid():
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if align_right:
            para.alignment = PP_ALIGN.RIGHT
        hl = blk.blockFormat().headingLevel()
        tl = blk.textList()
        if tl is not None:
            from PySide6.QtGui import QTextListFormat
            style = tl.format().style()
            ordered = style in (QTextListFormat.Style.ListDecimal,
                                QTextListFormat.Style.ListLowerAlpha,
                                QTextListFormat.Style.ListUpperAlpha,
                                QTextListFormat.Style.ListLowerRoman,
                                QTextListFormat.Style.ListUpperRoman)
            marker = (tl.itemText(blk) + " ") if ordered else "•  "
            _add_run(para, marker, base_pt, font_ref, color, bold=bold)
        it = blk.begin()
        while not it.atEnd():
            frag = it.fragment()
            if frag.isValid() and frag.text():
                _add_fragment(para, frag, base_pt, font_ref, color, bold, hl)
            it += 1
        blk = blk.next()


def _add_fragment(para, frag, base_pt: float, font_ref: str, color,
                  base_bold: bool, heading_level: int) -> None:
    cf = frag.charFormat()
    size = base_pt
    is_heading = heading_level > 0
    if is_heading:
        size = base_pt * _HEADING_SCALE.get(heading_level, 1.0)
    bold = base_bold or is_heading or cf.font().bold()
    href = cf.anchorHref() if cf.isAnchor() else None
    _add_run(para, frag.text(), size, font_ref, color, bold=bold,
             italic=cf.fontItalic(), code=cf.fontFixedPitch(), href=href)


def _add_run(para, text: str, size_pt: float, font_ref: str, color, *,
             bold: bool = False, italic: bool = False, code: bool = False,
             href: str | None = None) -> None:
    run = para.add_run()
    run.text = text
    font = run.font
    font.size = Pt(max(1.0, size_pt))
    font.bold = bold
    font.italic = italic
    font.name = "Consolas" if code else font_ref
    if href:
        run.hyperlink.address = href
        font.color.theme_color = MSO_THEME_COLOR.HYPERLINK
    else:
        font.color.theme_color = color


# ── raster ───────────────────────────────────────────────────────────────────

def _render_region(view, plan: SlidePlan, fitted: QRectF) -> QImage:
    """Rasterize the plan's framed scene region to a transparent QImage at the
    fitted size (capped), with overlay notes and container chrome hidden and an
    isolate scope applied — identical to the PDF exporter's diagram raster."""
    source = plan.source
    iw, ih = max(1, round(fitted.width())), max(1, round(fitted.height()))
    cap = 4096
    # Render at 4x slide points (~288 DPI, on par with the PDF's 300) so the
    # diagram stays crisp full-screen on hi-dpi displays; the cap keeps a
    # full-bleed hero bounded.
    scale = 4
    iw, ih = iw * scale, ih * scale
    if max(iw, ih) > cap:
        k = cap / max(iw, ih)
        iw, ih = max(1, round(iw * k)), max(1, round(ih * k))
    img = QImage(iw, ih, QImage.Format.Format_ARGB32_Premultiplied)
    img.fill(Qt.GlobalColor.transparent)
    ip = QPainter(img)
    ip.setRenderHint(QPainter.RenderHint.Antialiasing)
    ip.setRenderHint(QPainter.RenderHint.TextAntialiasing)
    ip.setRenderHint(QPainter.RenderHint.SmoothPixmapTransform)
    from grafli.pdfexport import _hidden
    with _hidden(list(plan.overlays) + list(plan.chrome_suppress)):
        if plan.isolate:
            with isolate_focus(view, plan.isolate):
                view._scene.render(ip, QRectF(0, 0, iw, ih), source)
        else:
            view._scene.render(ip, QRectF(0, 0, iw, ih), source)
    ip.end()
    return img


def _png_stream(img: QImage) -> io.BytesIO:
    ba = QByteArray()
    buf = QBuffer(ba)
    buf.open(QIODevice.OpenModeFlag.WriteOnly)
    img.save(buf, "PNG")
    buf.close()
    return io.BytesIO(bytes(ba))


def _footer_reserve(ctx: _Ctx) -> float:
    return ctx.page_h * _FOOTER_RESERVE_RATIO if ctx.footer else 0.0
