"""Tests for grafli.pptxexport — flow → editable PowerPoint.

These reopen the generated .pptx with python-pptx and assert structure (slide
count, native title/caption/footer text, embedded diagram pictures, theme colour
and font references, markdown run formatting). Because a .pptx is inspectable
XML, this verifies the real output rather than just that export ran.
"""

from __future__ import annotations

import os

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from PySide6.QtWidgets import QApplication

from grafli.format import parse
from grafli.pptxexport import export_flow_to_pptx

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")

# Three content stops: a normal diagram stop with a caption, a second diagram
# stop, and a text slide (single-note focus, no description). Plus a footer.
SAMPLE = """\
#!grafli v2
@ box a "A" 0,0 120x60
@ box b "B" 300,0 120x60
@ arrow a -> b "talks"
@ note ntxt 600,0 "md:\\n# Heading\\n\\n**bold** and a [link](https://example.com)\\n\\n- one\\n- two"
@ bookmark bm1 "Overview" @a,b "Both parts."
@ bookmark bm2 "Just A" @a "The left one."
@ bookmark bmtext "Notes" @ntxt
@ flow tour "Tour" bm1 bm2:5 bmtext "Wide then in."
@ footer "(c) Team"
"""


def _app() -> QApplication:
    return QApplication.instance() or QApplication([])


def _view(board):
    _app()
    from grafli.view import GrafliView
    v = GrafliView()
    v.load_board(board)
    return v


def _texts(slide):
    return [sh.text_frame.text for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]


def _has_picture(slide):
    return any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def _theme_root(prs):
    part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    return etree.fromstring(part.blob)


def _scheme_color(prs, slot):
    clr = _theme_root(prs).find(".//" + qn("a:clrScheme"))
    el = clr.find(qn("a:" + slot))[0]
    return el.get("val") or el.get("lastClr")


def test_export_produces_title_plus_one_slide_per_stop(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "tour.pptx"
    slides, _ = export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    assert slides == 4                       # title + 3 stops
    prs = Presentation(out)
    assert len(prs.slides._sldIdLst) == 4


def test_slide_size_is_16x9(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    # 960x540 pt in EMU (12700 EMU per pt).
    assert prs.slide_width == 960 * 12700
    assert prs.slide_height == 540 * 12700


def test_title_progress_and_caption_are_native_text(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    first = list(prs.slides)[1]              # first content stop (bm1)
    texts = _texts(first)
    assert "Overview" in texts
    assert "1 / 3" in texts
    assert any("Both parts." in t for t in texts)   # caption


def test_diagram_is_an_embedded_picture(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    assert _has_picture(list(prs.slides)[1])


def test_text_slide_is_native_text_without_picture(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    text_slide = list(prs.slides)[3]         # bmtext, single-note focus
    assert not _has_picture(text_slide)
    body = " ".join(_texts(text_slide))
    assert "bold" in body and "Heading" in body


def test_footer_present_on_content_slides(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    assert any("(c) Team" in t for t in _texts(list(prs.slides)[1]))


def test_markdown_bold_and_link_become_runs(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    text_slide = list(prs.slides)[3]
    runs = [r for sh in text_slide.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]
    assert any(r.font.bold for r in runs)                       # **bold**
    assert any(r.hyperlink.address for r in runs)               # [link](...)
    assert any(r.hyperlink.address == "https://example.com" for r in runs)


def test_runs_use_theme_font_and_colour_refs(tmp_path):
    # Theme-friendly: a global theme/font swap must cascade, so runs reference
    # the theme font (+mj-lt/+mn-lt) and a scheme colour rather than literals.
    from pptx.enum.dml import MSO_COLOR_TYPE
    board = parse(SAMPLE)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    title_run = None
    for sh in list(prs.slides)[1].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Overview":
            title_run = sh.text_frame.paragraphs[0].runs[0]
    assert title_run is not None
    assert title_run.font.name in ("+mj-lt", "+mn-lt")
    assert title_run.font.color.type == MSO_COLOR_TYPE.SCHEME


def test_grafli_theme_injects_palette_and_font(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "g.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                        theme="grafli")
    prs = Presentation(out)
    assert _scheme_color(prs, "accent1").upper() == "D4804E"   # grafli accent
    assert _scheme_color(prs, "lt1").upper() == "E8E4DD"       # paper bg
    major = _theme_root(prs).find(
        ".//" + qn("a:majorFont") + "/" + qn("a:latin")).get("typeface")
    assert "JetBrains" in major


def test_blank_theme_keeps_default_office_theme(tmp_path):
    board = parse(SAMPLE)
    out = tmp_path / "b.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                        theme="blank")
    prs = Presentation(out)
    assert _scheme_color(prs, "accent1").upper() != "D4804E"   # untouched


def test_graph_only_stop_has_picture_and_no_title(tmp_path):
    # Drop the footer so the graph-only stop is truly chrome-free (the footer is
    # board-global and shows on every content slide, including this one).
    src = SAMPLE.replace('@ footer "(c) Team"\n', '')
    src = src.replace('@ flow tour "Tour" bm1 bm2:5 bmtext',
                      '@ bookmark bm0 "" @a,b\n@ flow tour "Tour" bm0 bm1 bm2:5 bmtext')
    board = parse(src)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    prs = Presentation(out)
    graph_only = list(prs.slides)[1]         # bm0: no label, no description
    assert _has_picture(graph_only)
    assert _texts(graph_only) == []          # no title/progress/caption chrome


def test_container_box_label_titles_the_slide(tmp_path):
    src = """\
#!grafli v2
@ box frame "Frame" 0,0 400x200
@ box c1 "C" 200,40 120x60 >frame
@ note n1 20,20 "md:\\nhi" ~width=20 >frame
@ bookmark bmc "" @frame,n1,c1 ~iso
@ flow fc "FC" bmc
"""
    board = parse(src)
    out = tmp_path / "t.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("fc"), out)
    prs = Presentation(out)
    assert any(t == "Frame" for t in _texts(list(prs.slides)[1]))


def test_overload_reported_when_inplace_note_too_small(tmp_path):
    wide = """\
#!grafli v2
@ box frame "Frame" 0,0 4200x2400
@ box huge "Huge" 50,50 4000x2200 >frame
@ note cap 80,80 "md:\\nsmall caption text here" ~width=20 >huge
@ bookmark bmO "" @frame,huge,cap ~iso
@ flow f "F" bmO
"""
    board = parse(wide)
    out = tmp_path / "wide.pptx"
    _, overloaded = export_flow_to_pptx(_view(board), board.flow_by_id("f"), out)
    assert overloaded and overloaded[0][0] == 0


def test_export_restores_selection(tmp_path):
    board = parse(SAMPLE)
    view = _view(board)
    view._box_items["a"].setSelected(True)
    export_flow_to_pptx(view, board.flow_by_id("tour"), tmp_path / "t.pptx")
    assert view._box_items["a"].isSelected()


def _make_template(path, accent="ABCDEF"):
    """A minimal .pptx 'template': the default Office deck (4:3, with Title Slide
    and Title-and-Content layouts) but with a sentinel accent1 colour so we can
    prove the template theme survives the export, plus one pre-existing slide so
    we can prove slide-stripping works."""
    prs = Presentation()
    part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    th = etree.fromstring(part.blob)
    el = th.find(".//" + qn("a:clrScheme")).find(qn("a:accent1"))
    for c in list(el):
        el.remove(c)
    etree.SubElement(el, qn("a:srgbClr"), val=accent)
    part._blob = etree.tostring(th, xml_declaration=True, encoding="UTF-8",
                                standalone=True)
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return path


def _title_ph_text(slide):
    for ph in slide.placeholders:
        if "TITLE" in str(ph.placeholder_format.type):
            return ph.text_frame.text
    return None


def test_template_strips_existing_slides_and_keeps_step_count(tmp_path):
    tpl = _make_template(tmp_path / "tpl.pptx")
    board = parse(SAMPLE)
    out = tmp_path / "onto.pptx"
    slides, _ = export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                                    template=str(tpl))
    assert slides == 4                                  # title + 3 stops
    assert len(Presentation(out).slides._sldIdLst) == 4  # template slide gone


def test_template_preserves_its_theme_and_size(tmp_path):
    tpl = _make_template(tmp_path / "tpl.pptx", accent="ABCDEF")
    board = parse(SAMPLE)
    out = tmp_path / "onto.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                        template=str(tpl))
    prs = Presentation(out)
    assert _scheme_color(prs, "accent1").upper() == "ABCDEF"     # theme survives
    assert _scheme_color(prs, "accent1").upper() != "D4804E"     # not injected
    src = Presentation(tpl)
    assert prs.slide_width == src.slide_width                    # adopts size
    assert prs.slide_height == src.slide_height


def test_template_title_and_content_use_placeholders(tmp_path):
    tpl = _make_template(tmp_path / "tpl.pptx")
    board = parse(SAMPLE)
    out = tmp_path / "onto.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                        template=str(tpl))
    prs = Presentation(out)
    assert _title_ph_text(list(prs.slides)[0]) == "Tour"               # cover
    first = list(prs.slides)[1]                          # first stop (bm1)
    assert _title_ph_text(first) == "Overview"           # title placeholder
    assert _has_picture(first)                            # diagram embedded


def test_template_drops_grafli_progress_counter(tmp_path):
    # On a template the deck supplies its own chrome, so grafli's 'i / n' counter
    # is dropped (the template didn't budget space for it).
    tpl = _make_template(tmp_path / "tpl.pptx")
    board = parse(SAMPLE)
    out = tmp_path / "onto.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                        template=str(tpl))
    first = list(Presentation(out).slides)[1]
    assert not any("1 / 3" in t for t in _texts(first))


def test_template_explicit_layout_selection(tmp_path):
    tpl = _make_template(tmp_path / "tpl.pptx")
    board = parse(SAMPLE)
    out = tmp_path / "onto.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out,
                        template=str(tpl), title_layout="Title Slide",
                        content_layout="Title and Content")
    prs = Presentation(out)
    assert list(prs.slides)[0].slide_layout.name == "Title Slide"
    assert list(prs.slides)[1].slide_layout.name == "Title and Content"


def _note_run_sizes(slide, fragment):
    """Run sizes of the (single) shape containing ``fragment`` — scoped so the
    title bar / progress / footer chrome doesn't pollute the assertion."""
    for sh in slide.shapes:
        if sh.has_text_frame and fragment in sh.text_frame.text:
            return [r.font.size.pt for p in sh.text_frame.paragraphs
                    for r in p.runs if r.font.size is not None]
    return []


def test_playback_text_fit_parity_clamp_and_fallback():
    from PySide6.QtCore import QRectF
    from grafli.slideplan import playback_text_fit
    hero = QRectF(0, 0, 800, 400)
    # A 200x100 block fits the hero at 4x zoom: 16px text would read as 64 —
    # capped at the band max, the block itself filling the hero.
    size, rect = playback_text_fit(16, QRectF(0, 0, 200, 100), hero, 18, 60)
    assert size == 60
    assert (rect.width(), rect.height()) == (800, 400)
    # Aspect mismatch letterboxes: a wide block centres vertically.
    size, rect = playback_text_fit(16, QRectF(0, 0, 400, 100), hero, 18, 60)
    assert size == 32
    assert (rect.top(), rect.height()) == (100, 200)
    # A dense note lands below the band floor → None, caller band-fits.
    assert playback_text_fit(16, QRectF(0, 0, 200, 1000), hero, 18, 60) is None
    assert playback_text_fit(16, None, hero, 18, 60) is None


def test_short_text_slide_sized_at_playback_parity(tmp_path):
    # In the app, playback zooms a text step's note to fill the viewport; a
    # short note therefore reads big. The export mirrors that zoom (capped at
    # _BODY_MAX_PT) instead of typesetting it small mid-slide at the old band.
    src = SAMPLE.replace(
        '@ note ntxt 600,0 "md:\\n# Heading\\n\\n**bold** and a '
        '[link](https://example.com)\\n\\n- one\\n- two"',
        '@ note ntxt 600,0 "Ship it" ~40')
    board = parse(src)
    out = tmp_path / "deck.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    text_slide = list(Presentation(out).slides)[3]
    # 40px at ~2x playback zoom would read ~80 — clamped at _BODY_MAX_PT, far
    # above the old band cap of 30.
    assert _note_run_sizes(text_slide, "Ship it") == [60.0]


def test_dense_text_slide_falls_back_to_band(tmp_path):
    # A note too dense to be readable at its framed scale keeps the band's
    # shrink-to-fit instead of going sub-readable parity-small.
    lines = "\\n".join(f"line {i} of a long dense note body" for i in range(40))
    src = SAMPLE.replace(
        '@ note ntxt 600,0 "md:\\n# Heading\\n\\n**bold** and a '
        '[link](https://example.com)\\n\\n- one\\n- two"',
        f'@ note ntxt 600,0 "{lines}"')
    board = parse(src)
    out = tmp_path / "deck.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), out)
    text_slide = list(Presentation(out).slides)[3]
    sizes = _note_run_sizes(text_slide, "long dense note")
    assert sizes and all(18.0 <= s <= 30.0 for s in sizes)


def test_blank_theme_omits_title_accent_rule(tmp_path):
    # The decorative accent rule on the cover is grafli-only chrome; the blank
    # preset drops it (cleaner base for corporate templates), so its title slide
    # carries fewer shapes than grafli's for the same flow.
    board = parse(SAMPLE)
    g = tmp_path / "g.pptx"
    b = tmp_path / "b.pptx"
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), g, theme="grafli")
    export_flow_to_pptx(_view(board), board.flow_by_id("tour"), b, theme="blank")
    g_cover = len(list(Presentation(g).slides)[0].shapes)
    b_cover = len(list(Presentation(b).slides)[0].shapes)
    assert g_cover > b_cover
